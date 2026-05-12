"""
retrieval/document_loader.py
----------------------------
Memory-bounded document text extraction for PDF, DOCX, PPTX, TXT, and raw text.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass
from pathlib import Path
from typing import List, Union

logger = logging.getLogger(__name__)


@dataclass(frozen=True)
class DocumentLimits:
    max_document_chars: int = 2_000_000
    max_file_bytes: int = 25_000_000
    max_pdf_pages: int = 250
    max_pptx_slides: int = 250
    max_docx_paragraphs: int = 5000


class DocumentLoadError(RuntimeError):
    pass


def _trim_parts(parts: list[str], max_chars: int) -> list[str]:
    trimmed: list[str] = []
    total = 0
    for part in parts:
        text = str(part or "").strip()
        if not text:
            continue
        remaining = max_chars - total
        if remaining <= 0:
            break
        if len(text) > remaining:
            text = text[:remaining]
        trimmed.append(text)
        total += len(text)
    return trimmed


def _existing_path(source: Union[str, Path]) -> Path | None:
    if isinstance(source, Path):
        return source if source.exists() else None
    if not isinstance(source, str):
        return None
    if "\n" in source or "\r" in source or len(source) > 500:
        return None
    try:
        path = Path(source)
        return path if path.exists() else None
    except (OSError, ValueError):
        return None


def load_document(source: Union[str, Path], limits: DocumentLimits | None = None) -> List[str]:
    limits = limits or DocumentLimits()
    if isinstance(source, str) and not source.strip():
        return []
    path = _existing_path(source)

    if path is None:
        logger.info("[DOC] Treating input as raw text.")
        return _trim_parts([str(source)], limits.max_document_chars)

    if not path.is_file():
        raise DocumentLoadError(f"Document path is not a file: {path}")
    size = path.stat().st_size
    if size > limits.max_file_bytes:
        raise DocumentLoadError(
            f"File is too large ({size} bytes > limit {limits.max_file_bytes} bytes): {path.name}"
        )

    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _load_pdf(path, limits)
    if suffix == ".docx":
        return _load_docx(path, limits)
    if suffix in (".pptx", ".ppt"):
        return _load_pptx(path, limits)
    if suffix in (".txt", ".md", ".markdown", ".rst"):
        return _load_text(path, limits)

    logger.warning("[DOC] Unknown extension '%s'; trying plain text.", suffix)
    return _load_text(path, limits)


def _load_pdf(path: Path, limits: DocumentLimits) -> List[str]:
    try:
        import fitz

        doc = fitz.open(str(path))
        pages: list[str] = []
        try:
            for index, page in enumerate(doc):
                if index >= limits.max_pdf_pages:
                    break
                text = page.get_text("text").strip()
                if text:
                    pages.append(text)
        finally:
            doc.close()
        result = _trim_parts(pages, limits.max_document_chars)
        logger.info("[DOC] PDF '%s' -> %d page text block(s).", path.name, len(result))
        return result
    except ImportError:
        logger.warning("[DOC] PyMuPDF is unavailable; falling back to pypdf.")
    except Exception as exc:
        logger.warning("[DOC] PyMuPDF load failed for %s: %s; falling back to pypdf.", path.name, exc)

    try:
        import pypdf

        reader = pypdf.PdfReader(str(path), strict=False)
        pages: list[str] = []
        for index, page in enumerate(reader.pages):
            if index >= limits.max_pdf_pages:
                break
            text = (page.extract_text() or "").strip()
            if text:
                pages.append(text)
        result = _trim_parts(pages, limits.max_document_chars)
        logger.info("[DOC] PDF '%s' -> %d page text block(s).", path.name, len(result))
        return result
    except Exception as exc:
        logger.error("[DOC] PDF load failed for %s: %s", path.name, exc)
        raise DocumentLoadError(str(exc)) from exc


def _load_docx(path: Path, limits: DocumentLimits) -> List[str]:
    try:
        import docx

        doc = docx.Document(str(path))
        paragraphs = [
            p.text.strip()
            for p in doc.paragraphs[: limits.max_docx_paragraphs]
            if p.text and p.text.strip()
        ]
        result = _trim_parts(paragraphs, limits.max_document_chars)
        logger.info("[DOC] DOCX '%s' -> %d paragraph(s).", path.name, len(result))
        return result
    except Exception as exc:
        logger.error("[DOC] DOCX load failed for %s: %s", path.name, exc)
        raise DocumentLoadError(str(exc)) from exc


def _load_pptx(path: Path, limits: DocumentLimits) -> List[str]:
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        slides: list[str] = []
        for index, slide in enumerate(prs.slides):
            if index >= limits.max_pptx_slides:
                break
            texts = []
            for shape in slide.shapes:
                shape_text = getattr(shape, "text", "")
                if shape_text and shape_text.strip():
                    texts.append(shape_text.strip())
            if texts:
                slides.append("\n".join(texts))
        result = _trim_parts(slides, limits.max_document_chars)
        logger.info("[DOC] PPTX '%s' -> %d slide text block(s).", path.name, len(result))
        return result
    except Exception as exc:
        logger.error("[DOC] PPTX load failed for %s: %s", path.name, exc)
        raise DocumentLoadError(str(exc)) from exc


def _load_text(path: Path, limits: DocumentLimits) -> List[str]:
    try:
        text = path.read_text(encoding="utf-8", errors="replace").strip()
        result = _trim_parts([text], limits.max_document_chars)
        logger.info("[DOC] Text '%s' -> %d chars.", path.name, sum(len(p) for p in result))
        return result
    except Exception as exc:
        logger.error("[DOC] Text load failed for %s: %s", path.name, exc)
        raise DocumentLoadError(str(exc)) from exc
